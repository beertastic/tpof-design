#!/usr/bin/env python3
"""TPOF top-level character-board generator.

Examples:
  python tools/board-generator/generate.py shada
  python tools/board-generator/generate.py shada --board weapons
  python tools/board-generator/generate.py shada --dpi 150
  python tools/board-generator/generate.py shada --pdf-only
  python tools/board-generator/generate.py shada --validate
  python tools/board-generator/generate.py --all

The generator discovers character folders that already exist. It never creates
missing character directories when using --all, so deleted characters remain
deleted.
"""
from __future__ import annotations

import argparse
import copy
import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any

sys.path.insert(0, str(Path(__file__).parent))

import fitz
import yaml
import pdfrender
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DEFAULT_BOARD_ORDER = ["production", "costume", "weapons", "performance", "materials"]


def board_order(config: dict[str, Any]) -> list[str]:
    """Board order comes from the config; falls back to the classic five."""
    boards = config.get("boards") or {}
    return list(boards.keys()) if boards else list(DEFAULT_BOARD_ORDER)
A2_WIDTH_IN = 23.386
A2_HEIGHT_IN = 16.535


def find_repo_root(start: Path) -> Path:
    current = start.resolve()
    for candidate in [current, *current.parents]:
        if (candidate / "03-characters").is_dir() and (candidate / "tools").exists():
            return candidate
    raise SystemExit("Could not locate repository root containing 03-characters/ and tools/.")


def load_yaml(path: Path) -> dict[str, Any]:
    try:
        return yaml.safe_load(path.read_text(encoding="utf-8")) or {}
    except FileNotFoundError as exc:
        raise SystemExit(f"Missing configuration: {path}") from exc
    except yaml.YAMLError as exc:
        raise SystemExit(f"Invalid YAML in {path}: {exc}") from exc


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    if len(value) != 6:
        raise ValueError(f"Expected six-digit hex colour, got {hex_value!r}")
    return RGBColor(*(int(value[i:i+2], 16) for i in (0, 2, 4)))


def deep_merge(base: dict[str, Any], override: dict[str, Any]) -> dict[str, Any]:
    result = copy.deepcopy(base)
    for key, value in override.items():
        if isinstance(value, dict) and isinstance(result.get(key), dict):
            result[key] = deep_merge(result[key], value)
        else:
            result[key] = copy.deepcopy(value)
    return result


def resolve_character_dir(repo: Path, name: str) -> Path:
    direct = repo / "03-characters" / name
    if direct.is_dir():
        return direct
    lowered = name.casefold()
    matches = [p for p in (repo / "03-characters").iterdir() if p.is_dir() and p.name.casefold() == lowered]
    if len(matches) == 1:
        return matches[0]
    raise SystemExit(f"Character folder does not exist: 03-characters/{name}")


def discover_characters(repo: Path) -> list[Path]:
    found = []
    for path in sorted((repo / "03-characters").iterdir()):
        if path.is_dir() and (path / "board-data.yaml").is_file():
            found.append(path)
    return found


def validate(repo: Path, character_dir: Path, config: dict[str, Any], selected: set[str] | None = None) -> list[str]:
    errors: list[str] = []
    required = ["project", "character", "asset_id", "version", "status", "boards"]
    for key in required:
        if key not in config:
            errors.append(f"missing key: {key}")
    boards = config.get("boards", {})
    for board_name in board_order(config):
        if selected and board_name not in selected:
            continue
        board = boards.get(board_name)
        if not board:
            errors.append(f"missing board configuration: {board_name}")
            continue
        for item in board.get("images", []):
            rel = item.get("path")
            if not rel:
                errors.append(f"{board_name}: image entry missing path")
                continue
            image_path = character_dir / rel
            if not image_path.is_file():
                errors.append(f"{board_name}: missing image {rel}")
            else:
                try:
                    with Image.open(image_path) as im:
                        im.verify()
                except Exception as exc:
                    errors.append(f"{board_name}: unreadable image {rel}: {exc}")
        if not board.get("sections"):
            errors.append(f"{board_name}: no text sections")

        # Later elements paint over earlier ones, so overlaps silently hide content.
        boxes = ([(i.get("x"), i.get("y"), i.get("w"), i.get("h"), "image " +
                   str(i.get("path", "")).split("/")[-1])
                  for i in board.get("images", [])]
                 + [(sec.get("x"), sec.get("y"), sec.get("w"), sec.get("h"),
                     "panel " + str(sec.get("heading", "")))
                    for sec in board.get("sections", [])])
        for a in range(len(boxes)):
            for c in range(a + 1, len(boxes)):
                ax, ay, aw, ah, an = boxes[a]
                bx, by, bw, bh, bn = boxes[c]
                if None in (ax, ay, aw, ah, bx, by, bw, bh):
                    continue
                ox = min(ax + aw, bx + bw) - max(ax, bx)
                oy = min(ay + ah, by + bh) - max(ay, by)
                if ox > 0.02 and oy > 0.02:
                    errors.append(
                        f"{board_name}: {an} overlaps {bn} by "
                        f"{ox:.2f} x {oy:.2f} in — the later one will hide it")
    for rel in config.get("governing_documents", []):
        target = (character_dir / rel).resolve()
        if not target.exists():
            errors.append(f"missing governing document: {rel}")
    return errors


def add_rect(slide, x, y, w, h, fill, line):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line; shape.line.width = Pt(0.6)
    return shape


def add_text(slide, value, x, y, w, h, size, colour, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True
    paragraph = frame.paragraphs[0]; paragraph.text = str(value); paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = "Aptos"; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = colour
    return box


def add_image_contain(slide, path: Path, x, y, w, h, panel, line):
    add_rect(slide, x, y, w, h, panel, line)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    px, py = x + (w - dw) / 2, y + (h - dh) / 2
    slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(dw), Inches(dh))


def add_section(slide, section, style):
    x, y, w, h = (section[k] for k in ("x", "y", "w", "h"))
    add_rect(slide, x, y, w, h, style["panel"], style["line"])
    add_text(slide, section["heading"], x + 0.18, y + 0.12, w - 0.36, 0.34,
             section.get("heading_size", 11.5), style["accent"], True)
    box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.52), Inches(w - 0.34), Inches(h - 0.62))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True
    for index, item in enumerate(section.get("items", [])):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = "• " + str(item); p.space_after = Pt(section.get("space_after", 4))
        p.font.name = "Aptos"; p.font.size = Pt(section.get("font_size", 12)); p.font.color.rgb = style["ink"]


def make_deck(repo: Path, character_dir: Path, config: dict[str, Any], template: dict[str, Any], boards: list[str]) -> tuple[Path, list[str]]:
    style_cfg = deep_merge(template.get("style", {}), config.get("style", {}))
    style = {k: rgb(v) for k, v in style_cfg.items()}
    prs = Presentation(); prs.slide_width = Inches(A2_WIDTH_IN); prs.slide_height = Inches(A2_HEIGHT_IN)
    created = []
    for page_index, board_name in enumerate(boards, 1):
        board = config["boards"][board_name]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = style["background"]
        add_text(slide, config["character"].upper(), 0.45, 0.25, 2.4, 0.65, 30, style["ink"], True)
        add_text(slide, board["title"], 3.0, 0.30, 13.5, 0.45, 19, style["ink"], True)
        meta = f'{board.get("label", board_name.upper())} | ASSET {config["asset_id"]} | VERSION {config["version"]} | STATUS {config["status"]} | {config.get("date", "")}'
        add_text(slide, meta, 3.0, 0.78, 15.5, 0.25, 8.5, style["muted"])
        add_rect(slide, 0.45, 1.14, 22.45, 0.015, style["line"], style["line"])
        for image in board.get("images", []):
            add_image_contain(slide, character_dir / image["path"], image["x"], image["y"], image["w"], image["h"], style["image_panel"], style["line"])
        for section in board.get("sections", []):
            add_section(slide, section, style)
        add_text(slide, config.get("footer_left", "THE PRICE OF FREEDOM - PRODUCTION DESIGN BIBLE"), 0.45, 16.03, 7.6, 0.25, 8.2, style["muted"])
        add_text(slide, board.get("tagline", ""), 7.4, 16.03, 9.0, 0.25, 8.2, style["muted"], align=PP_ALIGN.CENTER)
        add_text(slide, f"PAGE {page_index} OF {len(boards)}", 20.3, 16.03, 2.6, 0.25, 8.2, style["muted"], align=PP_ALIGN.RIGHT)
        created.append(board_name)
    source_dir = character_dir / "source"; source_dir.mkdir(exist_ok=True)
    pptx_path = source_dir / f'{config["character"]}-Production-Boards.pptx'
    prs.save(pptx_path)
    return pptx_path, created


def convert_to_pdf(pptx_path: Path, output_dir: Path) -> Path:
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if not executable:
        raise SystemExit("LibreOffice is required for PDF export but was not found on PATH.")
    subprocess.run([executable, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(pptx_path)], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    pdf = output_dir / (pptx_path.stem + ".pdf")
    if not pdf.is_file():
        raise SystemExit(f"LibreOffice did not create expected PDF: {pdf}")
    return pdf


def split_and_render(character_dir: Path, config: dict[str, Any], master_pdf: Path, boards: list[str], dpi: int, pdf_only: bool):
    doc = fitz.open(master_pdf)
    renders = character_dir / "renders"; renders.mkdir(exist_ok=True)
    outputs = []
    for index, board_name in enumerate(boards):
        board = config["boards"][board_name]
        default_pdf = board_name.replace("_", "-").title() + "-Board.pdf"
        output_pdf = character_dir / board.get("output_pdf", default_pdf)
        single = fitz.open(); single.insert_pdf(doc, from_page=index, to_page=index); single.save(output_pdf); single.close()
        outputs.append(output_pdf)
        if not pdf_only:
            opened = fitz.open(output_pdf)
            pix = opened[0].get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72), alpha=False)
            png_name = board.get("output_png") if dpi == 300 else None
            png_name = png_name or output_pdf.stem + f"-A2-{dpi}dpi.png"
            png_path = renders / png_name
            pix.save(png_path); opened.close(); outputs.append(png_path)
    doc.close()
    return outputs


def run_character(repo: Path, character_dir: Path, args, template: dict[str, Any]):
    config = load_yaml(character_dir / "board-data.yaml")
    selected = [args.board] if args.board else board_order(config)
    missing = [name for name in selected if name not in config.get("boards", {})]
    if missing:
        raise SystemExit(f"Missing configured board(s) for {character_dir.name}: {', '.join(missing)}")
    errors = validate(repo, character_dir, config, set(selected))
    if errors:
        print(f"Validation failed for {character_dir.relative_to(repo)}:", file=sys.stderr)
        for error in errors: print(f"  - {error}", file=sys.stderr)
        raise SystemExit(2)
    print(f"Validated {character_dir.relative_to(repo)}")
    if args.validate:
        return
    style_cfg = deep_merge(template.get("style", {}), config.get("style", {}))
    style = {k: str(v) for k, v in style_cfg.items()}

    outputs = []
    renders = character_dir / "renders"
    for index, board_name in enumerate(selected, 1):
        board = config["boards"][board_name]
        default_pdf = board_name.replace("_", "-").title() + "-Board.pdf"
        out_pdf = character_dir / board.get("output_pdf", default_pdf)
        pdfrender.render_board(out_pdf, character_dir, config, style, board,
                               index, len(selected), A2_WIDTH_IN, A2_HEIGHT_IN)
        outputs.append(out_pdf)

        if not args.pdf_only:
            renders.mkdir(exist_ok=True)
            doc = fitz.open(out_pdf)
            pix = doc[0].get_pixmap(matrix=fitz.Matrix(args.dpi / 72, args.dpi / 72),
                                    alpha=False)
            name = board.get("output_png") if args.dpi == 300 else None
            name = name or out_pdf.stem + f"-A2-{args.dpi}dpi.png"
            png = renders / name
            pix.save(png)
            doc.close()
            outputs.append(png)

    print("Generated:")
    for output in outputs:
        print(f"  {output.relative_to(repo)}")


def main():
    parser = argparse.ArgumentParser(description="Generate TPOF A2 character production boards.")
    parser.add_argument("character", nargs="?", help="Existing character directory name, e.g. shada")
    parser.add_argument("--all", action="store_true", help="Generate every existing character folder containing board-data.yaml")
    parser.add_argument("--board", help="Generate one board only, by key")
    parser.add_argument("--dpi", type=int, default=300, help="PNG render DPI (default: 300)")
    parser.add_argument("--pdf-only", action="store_true", help="Skip PNG preview rendering")
    parser.add_argument("--validate", action="store_true", help="Validate configuration and artwork without generating files")
    args = parser.parse_args()
    if args.all == bool(args.character):
        parser.error("Provide exactly one character name or use --all.")
    if args.dpi < 72 or args.dpi > 600:
        parser.error("--dpi must be between 72 and 600.")
    repo = find_repo_root(Path(__file__).parent)
    template = load_yaml(Path(__file__).parent / "templates" / "character-a2.yaml")
    characters = discover_characters(repo) if args.all else [resolve_character_dir(repo, args.character)]
    if not characters:
        raise SystemExit("No existing character folders contain board-data.yaml.")
    for character_dir in characters:
        run_character(repo, character_dir, args, template)

if __name__ == "__main__":
    main()
